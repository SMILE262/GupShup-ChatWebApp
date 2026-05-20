from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

slides = [
    {
        'title': 'GupShup Chat App',
        'content': [
            'A modern chat application built with React, Node.js, Express, MongoDB, and Socket.IO.',
            'Real-time messaging, authenticated user sessions, and conversation storage.'
        ]
    },
    {
        'title': 'Project Overview',
        'content': [
            'Client: React + Vite frontend for chat UI and theme selection.',
            'Server: Express API with MongoDB for users, conversations, and messages.',
            'Socket.IO for live online users and message delivery.'
        ]
    },
    {
        'title': 'Key Features',
        'content': [
            'Send and receive messages in real time.',
            'Store conversations and message history in MongoDB.',
            'Show online users with Socket.IO tracking.',
            'Theme and background selection in UI.'
        ]
    },
    {
        'title': 'Frontend Architecture',
        'content': [
            'React pages under client/src/pages/home.',
            'Redux Toolkit state management for users and messages.',
            'Axios instance for authenticated API calls.',
            'Dynamic UI components for chat bubbles and message input.'
        ]
    },
    {
        'title': 'Backend Architecture',
        'content': [
            'Express server with APIs in server/routes and server/controller.',
            'Middleware for authentication, error handling, and async operations.',
            'Socket server in server/socket/socket.js to manage online users.',
            'MongoDB models for User, Conversation, and Message.'
        ]
    },
    {
        'title': 'Data Model',
        'content': [
            'Message document stores senderId, receiverId, message text, timestamps.',
            'Conversation document stores participant user IDs and message references.',
            'Messages are populated from conversation when chat loads.'
        ]
    },
    {
        'title': 'Message Flow',
        'content': [
            'Client sends POST /message/send/:receiverId.',
            'Server creates Message and updates Conversation.',
            'If receiver is online, Socket.IO emits newMessage.',
            'Client fetches chat history with GET /message/getmessages/:ChatPartnerId.'
        ]
    },
    {
        'title': 'Delete Message Support',
        'content': [
            'Added DELETE /message/:messageId route.',
            'Backend removes message and pulls it from the conversation.',
            'Frontend uses deleteMessageThunk to update Redux state.',
            'UI shows delete options on double-click or hold.'
        ]
    },
    {
        'title': 'Authentication & Security',
        'content': [
            'Protected routes via isAuthenticated middleware.',
            'JWT-based authentication likely used in server auth.',
            'Users can only delete messages from their own conversations.'
        ]
    },
    {
        'title': 'User Experience',
        'content': [
            'Interactive chat UI with avatars and timestamps.',
            'Theme and background customization saved in localStorage.',
            'Mobile-friendly interactions with double-tap/hold delete.'
        ]
    },
    {
        'title': 'Project Tech Stack',
        'content': [
            'Frontend: React, Vite, Redux Toolkit, Tailwind-style classes.',
            'Backend: Node.js, Express, MongoDB, Mongoose, Socket.IO.',
            'Utilities: Axios, JSON Web Tokens, bcryptjs, cors, dotenv.'
        ]
    },
    {
        'title': 'Future Improvements',
        'content': [
            'Add chat clear or full conversation delete feature.',
            'Implement message edit and read receipts.',
            'Add notifications and group chat support.',
            'Improve visual design and accessibility.'
        ]
    }
]

for slide_data in slides:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_data['title']
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for idx, line in enumerate(slide_data['content']):
        if idx == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

prs.save(r'c:\Users\kumar\Documents\GupShup\GupShup_Presentation.pptx')
